import os
import argparse
import copy
import six
import datetime

from pptx import Presentation
from pptx.dml.color import ColorFormat, RGBColor

blank_slide_layout = 6

def get_arguments():
    parser = argparse.ArgumentParser()
    parser.add_argument('--path', type=str, default='none', help='path to pptx')
    args = parser.parse_args()
    return args

def main():
    args = get_arguments()
    if args.path == 'none':
        exit('please enter --path for path to pptx')

    dirs = os.listdir(args.path)
    prs = Presentation()
    for file in dirs:
        if file.endswith('.pptx') and not file.startswith('~'):
            print(file)
            temp_prs = Presentation(args.path + '/' + file)
            print('total: {}'.format(len(temp_prs.slides)))
            for idx in range(len(temp_prs.slides)):
                slide = temp_prs.slides[idx]
                # print(idx, slide, temp_prs.slide_layouts[4])
                new_slide = prs.slides.add_slide(temp_prs.slide_layouts[4])

                for shp in slide.shapes:
                    el = shp.element
                    newel = copy.deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                '''
                for _, value in six.iteritems(slide.part.rels):
                    # Make sure we don't copy a notesSlide relation as that won't exist
                    if "notesSlide" not in value.reltype:
                        new_slide.part.rels.add_relationship(value.reltype, value._target, value.rId)
                '''

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0)

    prs.save('a.pptx')
    prs.save(os.path.join(args.path, 'concated{}.pptx'.format(datetime.datetime.now().strftime("%Y%m%d%H%M%S"))))


if __name__ == '__main__':
    main()